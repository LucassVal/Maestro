#!/usr/bin/env python
# -*- coding: utf-8 -*-
import os
import sys
from .extrator_base_SRC_GL import ExtratorBase

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

class ExtratorPPTX(ExtratorBase):
    def __init__(self):
        super().__init__(nome="ExtratorPPTX")

    def executar(self, entrada):
        caminho = entrada.get("caminho")
        valido, msg = self.validar_arquivo(caminho)
        if not valido:
            return {"status": "ERROR", "error": msg}
        
        if not Presentation:
            return {"status": "ERROR", "error": "Biblioteca python-pptx não instalada."}

        try:
            prs = Presentation(caminho)
            texto = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texto += shape.text + "\n"
            return {
                "status": "SUCCESS",
                "conteudo": texto,
                "metadados": {"slides": len(prs.slides), "formato": "PPTX"}
            }
        except Exception as e:
            return {"status": "ERROR", "error": str(e)}

def executar(entrada):
    return ExtratorPPTX().executar(entrada)
